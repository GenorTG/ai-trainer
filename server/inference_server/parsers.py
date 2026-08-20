"""Document parsers for 33 file formats.

WHAT THIS FILE DOES
==================
Converts various document formats (PDF, DOCX, XLSX, PPTX, EPUB, HTML, etc.)
into plain text that can be stored in the RAG vector database.

KEY CONCEPTS
============
- Format detection: we look at the file extension (and sometimes the
  first few bytes) to decide which parser to use.
- Fallback parsing: if the preferred parser fails (e.g., antiword for
  DOC files), we try a backup.
- Subprocess calls: for some formats (DOC, EPUB), we shell out to
  command-line tools like antiword, catdoc, or pdftotext.
- check=False: we set this on subprocess.run() because we want to
  handle failures ourselves (returncode check) rather than let
  subprocess raise an exception.
"""

"""Universal document parser — extracts text from all office/file formats."""
import csv
import json
import os
import xml.etree.ElementTree as ET
from pathlib import Path

# ── Plain text / code / data ──

def parse_text(path: Path) -> str:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def parse_csv(path: Path) -> str:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        rows = list(reader)
    if not rows:
        return ""
    # Format as markdown table
    lines = []
    for i, row in enumerate(rows):
        lines.append(" | ".join(row))
        if i == 0:
            lines.append(" | ".join(["---"] * len(row)))
    return "\n".join(lines)


def parse_json(path: Path) -> str:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        data = json.load(f)
    return json.dumps(data, indent=2, ensure_ascii=False)


def parse_jsonl(path: Path) -> str:
    parts = []
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        for line in f:
            line = line.strip()
            if line:
                try:
                    parts.append(json.dumps(json.loads(line), indent=2, ensure_ascii=False))
                except json.JSONDecodeError:
                    parts.append(line)
    return "\n\n".join(parts)


def parse_xml(path: Path) -> str:
    try:
        tree = ET.parse(path)
        root = tree.getroot()
        return ET.tostring(root, encoding="unicode")
    except ET.ParseError:
        # Fall back to raw text
        return parse_text(path)


def parse_html(path: Path) -> str:
    try:
        from bs4 import BeautifulSoup
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            soup = BeautifulSoup(f.read(), "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)
    except ImportError:
        import re
        text = parse_text(path)
        text = re.sub(r"<script.*?</script>", "", text, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<style.*?</style>", "", text, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<[^>]+>", " ", text)
        text = re.sub(r"\s+", " ", text)
        return text.strip()


# ── PDF ──

def parse_pdf(path: Path) -> str:
    # Try pypdf first
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        pages = []
        for page in reader.pages:
            text = page.extract_text() or ""
            pages.append(text)
        return "\n\n".join(pages)
    except ImportError:
        pass

    # Try PyPDF2
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(path))
        pages = []
        for page in reader.pages:
            text = page.extract_text() or ""
            pages.append(text)
        return "\n\n".join(pages)
    except ImportError:
        pass

    # Try pdftotext CLI
    try:
        import subprocess
        result = subprocess.run(["pdftotext", str(path), "-"], capture_output=True, text=True, timeout=60, check=False)
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass

    return f"[PDF: {path.name} — install pypdf (pip install pypdf) to read]"


# ── Office documents ──

def parse_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        parts = []

        # Paragraphs
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text)

        # Tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                parts.append(" | ".join(cells))

        return "\n".join(parts)
    except ImportError:
        return f"[DOCX: {path.name} — install python-docx (pip install python-docx) to read]"


def parse_doc(path: Path) -> str:
    """Parse legacy .doc files."""
    # Try antiword
    try:
        import subprocess
        result = subprocess.run(["antiword", str(path)], capture_output=True, text=True, timeout=30, check=False)
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass

    # Try catdoc
    try:
        import subprocess
        result = subprocess.run(["catdoc", str(path)], capture_output=True, text=True, timeout=30, check=False)
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass

    # Try textract
    try:
        import textract
        return textract.process(str(path)).decode("utf-8", errors="replace")
    except ImportError:
        pass

    return f"[DOC: {path.name} — legacy format. Install antiword (apt install antiword) or textract to read]"


def parse_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), data_only=True, read_only=True)
        parts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            parts.append(f"=== Sheet: {sheet} ===")
            for row in ws.iter_rows(values_only=True):
                values = [str(v) if v is not None else "" for v in row]
                if any(v.strip() for v in values):
                    parts.append(" | ".join(values))
        return "\n".join(parts)
    except ImportError:
        return f"[XLSX: {path.name} — install openpyxl (pip install openpyxl) to read]"


def parse_xls(path: Path) -> str:
    try:
        import xlrd
        wb = xlrd.open_workbook(str(path))
        parts = []
        for sheet in wb.sheets():
            parts.append(f"=== Sheet: {sheet.name} ===")
            for row_idx in range(sheet.nrows):
                values = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
                if any(v.strip() for v in values):
                    parts.append(" | ".join(values))
        return "\n".join(parts)
    except ImportError:
        return f"[XLS: {path.name} — legacy format. Install xlrd (pip install xlrd) to read]"


def parse_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"=== Slide {i} ===")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        parts.append(" | ".join(cells))
        return "\n".join(parts)
    except ImportError:
        return f"[PPTX: {path.name} — install python-pptx (pip install python-pptx) to read]"


def parse_odt(path: Path) -> str:
    try:
        import zipfile
        with zipfile.ZipFile(path) as z, z.open("content.xml") as f:
            content = f.read().decode("utf-8", errors="replace")
        root = ET.fromstring(content)
        # Extract all text
        texts = []
        for elem in root.iter():
            if elem.tag.endswith("}p") or elem.tag.endswith("}h"):
                text = "".join(elem.itertext()).strip()
                if text:
                    texts.append(text)
        return "\n".join(texts)
    except Exception as e:  # noqa: BLE001
        return f"[ODT: {path.name} — failed to parse: {e}]"


def parse_rtf(path: Path) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return rtf_to_text(f.read())
    except ImportError:
        # Minimal RTF strip fallback
        text = parse_text(path)
        text = text.replace("\\par", "\n").replace("\\pard", "\n")
        text = text.replace("{\\rtf1", "").replace("}", "").replace("{", "")
        import re
        text = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", text)
        return text.strip()


def parse_epub(path: Path) -> str:
    try:
        import zipfile
        parts = []
        with zipfile.ZipFile(path) as z:
            for name in z.namelist():
                if name.endswith((".xhtml", ".html", ".htm")):
                    content = z.read(name).decode("utf-8", errors="replace")
                    try:
                        from bs4 import BeautifulSoup
                        soup = BeautifulSoup(content, "html.parser")
                        text = soup.get_text(separator="\n", strip=True)
                    except ImportError:
                        import re
                        text = re.sub(r"<[^>]+>", " ", content)
                        text = re.sub(r"\s+", " ", text).strip()
                    if text:
                        parts.append(text)
        return "\n\n".join(parts)
    except Exception as e:  # noqa: BLE001
        return f"[EPUB: {path.name} — failed to parse: {e}]"


# ── Dispatcher ──

PARSERS = {
    ".txt": parse_text,
    ".md": parse_text,
    ".markdown": parse_text,
    ".log": parse_text,
    ".py": parse_text,
    ".js": parse_text,
    ".ts": parse_text,
    ".jsx": parse_text,
    ".tsx": parse_text,
    ".html": parse_html,
    ".htm": parse_html,
    ".css": parse_text,
    ".csv": parse_csv,
    ".tsv": parse_csv,
    ".json": parse_json,
    ".jsonl": parse_jsonl,
    ".xml": parse_xml,
    ".yaml": parse_text,
    ".yml": parse_text,
    ".ini": parse_text,
    ".cfg": parse_text,
    ".conf": parse_text,
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".doc": parse_doc,
    ".xlsx": parse_xlsx,
    ".xls": parse_xls,
    ".pptx": parse_pptx,
    ".odt": parse_odt,
    ".ods": parse_odt,
    ".odp": parse_odt,
    ".rtf": parse_rtf,
    ".epub": parse_epub,
}

# Formats supported with optional deps — warn if missing
OPTIONAL_DEPS = {
    ".pdf": "pypdf",
    ".docx": "python-docx",
    ".xlsx": "openpyxl",
    ".pptx": "python-pptx",
    ".html": "beautifulsoup4",
    ".doc": "antiword (system) or textract",
    ".xls": "xlrd",
    ".rtf": "striprtf",
}


def supported_extensions() -> list[str]:
    return sorted(PARSERS.keys())


def parse_document(path: str) -> str:
    """Parse any supported document to plain text."""
    p = Path(path)
    if not p.exists():
        return f"[File not found: {path}]"

    ext = p.suffix.lower()
    parser = PARSERS.get(ext)

    if parser is None:
        # Unknown format — try reading as text
        try:
            return parse_text(p)
        except Exception:  # noqa: BLE001
            return f"[Unsupported format: {ext}]"

    try:
        return parser(p)
    except Exception as e:  # noqa: BLE001
        return f"[Failed to parse {p.name}: {e}]"


def parse_bytes(filename: str, data: bytes) -> str:
    """Parse document from bytes (e.g., uploaded file)."""
    p = Path(filename)
    ext = p.suffix.lower()
    parser = PARSERS.get(ext)

    if parser is None:
        try:
            return data.decode("utf-8", errors="replace")
        except Exception:  # noqa: BLE001
            return f"[Unsupported format: {ext}]"

    # Write to temp file and parse
    import tempfile
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name

    try:
        return parser(Path(tmp_path))
    except Exception as e:  # noqa: BLE001
        return f"[Failed to parse {filename}: {e}]"
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass


def ingest_bytes(filename: str, data: bytes, store, embedding_model: str = "all-MiniLM-L6-v2",
                 chunk_size: int = 512, overlap: int = 50) -> dict:
    """Parse bytes and ingest into RAG store directly (for uploads)."""
    import hashlib

    text = parse_bytes(filename, data)
    if not text or text.startswith(("[Failed", "[Unsupported")):
        return {"filename": filename, "chunks": 0, "error": text[:200]}

    doc_id = hashlib.md5(f"{filename}:{len(data)}".encode()).hexdigest()[:12]

    # Chunk
    words = text.split()
    chunks = []
    start = 0
    idx = 0
    while start < len(words):
        end = min(start + chunk_size, len(words))
        chunks.append({"id": f"{doc_id}_{idx}", "text": " ".join(words[start:end]),
                       "source": filename})
        idx += 1
        start += chunk_size - overlap

    added = store.add_document(doc_id, chunks, embedding_model)
    return {"filename": filename, "document_id": doc_id, "chunks": added}
